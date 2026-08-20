"""Zero-LLM PowerPoint extraction for OKF/RAG business knowledge."""
from __future__ import annotations

import io
import re
import uuid
from datetime import datetime, timezone

from features.okf_knowledge.md_extractor import (
    _clean_text,
    _split_oversized,
)

try:
    from pptx import Presentation

    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False


def pptx_extraction_available() -> bool:
    return _PPTX_AVAILABLE


def _document_metadata(source_filename: str) -> dict:
    match = re.search(
        r"\b(INS-(?:CFO|FIN|RESULTS)-[A-Z0-9-]+)\b",
        source_filename,
        re.IGNORECASE,
    )
    return {
        "doc_code": match.group(1).upper() if match else "",
        "doc_type": "quarterly_results",
        "owner": "Chief Financial Officer",
    }


def _shape_text(shape) -> list[str]:
    values: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = _clean_text(shape.text_frame.text or "")
        if text:
            values.append(text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [_clean_text(cell.text or "") for cell in row.cells]
            if any(cells):
                values.append(" | ".join(cells))
    return values


def extract_pptx_to_concepts(
    file_bytes: bytes,
    source_filename: str,
    default_tags: list[str] | None = None,
) -> list[dict]:
    """Extract slide text and tables into cited, slide-level concepts."""
    if not _PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is required for PowerPoint extraction.")

    presentation = Presentation(io.BytesIO(file_bytes))
    doc_id = str(uuid.uuid4())[:8]
    ingested_at = datetime.now(timezone.utc).isoformat()
    metadata = _document_metadata(source_filename)
    base_tags = list(
        default_tags
        or ["powerpoint", "business-doc", "insurance", "quarterly-results", "cfo"]
    )
    concepts: list[dict] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        blocks: list[str] = []
        for shape in slide.shapes:
            blocks.extend(_shape_text(shape))
        cleaned = _clean_text("\n\n".join(blocks))
        if len(cleaned) < 40:
            continue

        title = ""
        try:
            title = _clean_text(slide.shapes.title.text or "")
        except Exception:
            title = ""
        title = title or f"{source_filename} — slide {slide_number}"

        for chunk in _split_oversized(cleaned):
            body = chunk if chunk.startswith(title) else f"{title}\n\n{chunk}"
            concepts.append(
                {
                    "concept_id": f"{doc_id}_{slide_number}_{len(concepts)}",
                    "title": title[:90],
                    "source_doc": source_filename,
                    "source_page": slide_number,
                    "source_locator": f"slide {slide_number}",
                    "tags": base_tags,
                    "body": body[:1400],
                    "ingested_at": ingested_at,
                    **metadata,
                }
            )

    return concepts


def extract_pptx_file(path) -> list[dict]:
    from pathlib import Path

    file_path = Path(path)
    return extract_pptx_to_concepts(file_path.read_bytes(), file_path.name)
