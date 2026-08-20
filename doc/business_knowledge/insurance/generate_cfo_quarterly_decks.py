"""Generate two synthetic insurance CFO decks used by the ASK-DB RAG demo."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(14, 165, 233)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
WHITE = RGBColor(248, 250, 252)
MUTED = RGBColor(148, 163, 184)


Q1 = {
    "quarter": "Q1 2026",
    "code": "INS-CFO-Q1-2026",
    "filename": "INS-CFO-Q1-2026_Quarterly_Insurance_Results.pptx",
    "gwp": "₹1.24B",
    "earned": "₹1.10B",
    "incurred": "₹682M",
    "paid": "₹510M",
    "claims": "18,400",
    "loss_ratio": "62.0%",
    "severity": "₹37.1K",
    "approval": "83.4%",
    "settlement": "17.8 days",
    "renewal": "78.2%",
    "expense_ratio": "28.1%",
    "combined_ratio": "90.1%",
    "capital": "1.82×",
    "cash": "₹2.08B",
    "summary": [
        "Gross written premium reached ₹1.24B, 6.0% above Q4 2025.",
        "Loss ratio closed at 62.0%; Motor remained the largest source of volatility.",
        "Combined ratio was 90.1%, within the CFO operating threshold of below 95%.",
        "Renewal rate improved to 78.2%, supported by digital reminders.",
    ],
    "lob": [
        ["Line of business", "GWP", "Earned premium", "Loss ratio", "CFO comment"],
        ["Motor", "₹610M", "₹540M", "68%", "Repair inflation remains elevated"],
        ["Health", "₹390M", "₹350M", "58%", "Stable utilization and pricing"],
        ["Property", "₹240M", "₹210M", "49%", "Benign catastrophe experience"],
    ],
    "regions": [
        ["Region", "GWP", "Loss ratio", "Claim cycle", "Status"],
        ["West", "₹360M", "66%", "18.4 days", "Watch"],
        ["North", "₹335M", "64%", "19.1 days", "Watch"],
        ["South", "₹315M", "58%", "15.8 days", "On plan"],
        ["East", "₹230M", "57%", "17.2 days", "On plan"],
    ],
    "drivers": [
        "Motor own-damage severity increased due to parts and labour inflation.",
        "Digital first-notice-of-loss reduced simple-claim handling time by 1.6 days.",
        "No material catastrophe event was recorded in the quarter.",
        "Reserve review found no material adverse development.",
    ],
    "actions": [
        "Renegotiate preferred repair-network rates before the Q2 renewal cycle.",
        "Expand straight-through processing for low-value motor claims.",
        "Maintain pricing discipline in West and North motor portfolios.",
        "Target a Q2 combined ratio below 92% and settlement time below 18 days.",
    ],
}

Q2 = {
    "quarter": "Q2 2026",
    "code": "INS-CFO-Q2-2026",
    "filename": "INS-CFO-Q2-2026_Quarterly_Insurance_Results.pptx",
    "gwp": "₹1.33B",
    "earned": "₹1.18B",
    "incurred": "₹767M",
    "paid": "₹588M",
    "claims": "20,600",
    "loss_ratio": "65.0%",
    "severity": "₹37.2K",
    "approval": "81.7%",
    "settlement": "19.6 days",
    "renewal": "76.5%",
    "expense_ratio": "27.5%",
    "combined_ratio": "92.5%",
    "capital": "1.76×",
    "cash": "₹2.16B",
    "summary": [
        "Gross written premium grew 7.3% quarter on quarter to ₹1.33B.",
        "Loss ratio deteriorated 3.0 points to 65.0%, driven by Motor and weather losses.",
        "Combined ratio increased to 92.5% but remained below the 95% CFO threshold.",
        "Claims settlement time rose to 19.6 days; remediation is now an executive priority.",
    ],
    "lob": [
        ["Line of business", "GWP", "Earned premium", "Loss ratio", "CFO comment"],
        ["Motor", "₹665M", "₹585M", "74%", "Weather and repair inflation"],
        ["Health", "₹410M", "₹365M", "61%", "Higher elective utilization"],
        ["Property", "₹255M", "₹230M", "47%", "Pricing and risk selection improved"],
    ],
    "regions": [
        ["Region", "GWP", "Loss ratio", "Claim cycle", "Status"],
        ["West", "₹390M", "72%", "21.3 days", "Escalate"],
        ["North", "₹355M", "69%", "20.8 days", "Escalate"],
        ["South", "₹340M", "59%", "17.6 days", "Watch"],
        ["East", "₹245M", "58%", "18.2 days", "Watch"],
    ],
    "drivers": [
        "Motor weather-related claims added approximately 2.1 points to the loss ratio.",
        "Repair invoice inflation continued, especially in West and North.",
        "Approval rate declined as fraud-review referrals increased.",
        "Property underwriting actions reduced loss ratio by 2 points quarter on quarter.",
    ],
    "actions": [
        "Launch a 60-day Motor loss-ratio recovery program in West and North.",
        "Add adjuster capacity and daily SLA escalation for claims older than 15 days.",
        "Apply targeted Motor pricing action at renewal; avoid broad portfolio discounting.",
        "Target Q3 loss ratio below 63%, settlement below 18 days, and renewal above 77%.",
    ],
}


def _set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY


def _add_text(slide, text, left, top, width, height, size=18, colour=WHITE, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = colour
    return box


def _add_header(slide, title, subtitle=""):
    _set_background(slide)
    _add_text(slide, title, 0.65, 0.35, 11.8, 0.55, 26, WHITE, True)
    if subtitle:
        _add_text(slide, subtitle, 0.68, 0.9, 11.7, 0.35, 11, MUTED)
    line = slide.shapes.add_shape(1, Inches(0.65), Inches(1.25), Inches(12), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = CYAN
    line.line.fill.background()


def _add_footer(slide, data):
    _add_text(
        slide,
        f"{data['code']} · SYNTHETIC DEMO DATA — NOT A REAL COMPANY DISCLOSURE",
        0.65,
        7.08,
        11.9,
        0.22,
        8,
        MUTED,
    )


def _add_bullets(slide, bullets, left=0.85, top=1.65, width=11.5, height=4.8):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    for index, item in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = WHITE
        paragraph.space_after = Pt(16)
        paragraph.text = f"• {item}"


def _add_table(slide, rows, left=0.7, top=1.55, width=11.95, height=4.9):
    table = slide.shapes.add_table(
        len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height)
    ).table
    for row_index, row in enumerate(rows):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if row_index == 0 else RGBColor(30, 41, 59)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Aptos"
                paragraph.font.size = Pt(11 if row_index else 12)
                paragraph.font.bold = row_index == 0
                paragraph.font.color.rgb = WHITE
                paragraph.alignment = PP_ALIGN.LEFT
    return table


def _add_kpi_grid(slide, data):
    kpis = [
        ("GWP", data["gwp"], BLUE),
        ("Earned premium", data["earned"], CYAN),
        ("Claims incurred", data["incurred"], AMBER),
        ("Claims paid", data["paid"], AMBER),
        ("Claim count", data["claims"], BLUE),
        ("Loss ratio", data["loss_ratio"], RED if data["quarter"].startswith("Q2") else AMBER),
        ("Combined ratio", data["combined_ratio"], GREEN),
        ("Renewal rate", data["renewal"], CYAN),
    ]
    for index, (label, value, colour) in enumerate(kpis):
        row, col = divmod(index, 4)
        left = 0.7 + col * 3.05
        top = 1.55 + row * 2.3
        shape = slide.shapes.add_shape(
            5, Inches(left), Inches(top), Inches(2.75), Inches(1.85)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(30, 41, 59)
        shape.line.color.rgb = colour
        _add_text(slide, label, left + 0.18, top + 0.2, 2.4, 0.3, 11, MUTED)
        _add_text(slide, value, left + 0.18, top + 0.62, 2.4, 0.6, 25, colour, True)


def build_deck(data):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    _set_background(slide)
    _add_text(slide, "ASK-DB Insurance", 0.8, 0.7, 5.8, 0.45, 17, CYAN, True)
    _add_text(slide, f"CFO Quarterly Results · {data['quarter']}", 0.8, 1.55, 11.5, 0.9, 34, WHITE, True)
    _add_text(slide, "Financial performance, claims, capital, outlook and actions", 0.82, 2.55, 10.8, 0.45, 18, MUTED)
    _add_text(slide, data["code"], 0.82, 3.3, 4.5, 0.35, 14, CYAN, True)
    _add_text(slide, "SYNTHETIC DEMO DATA", 0.82, 4.25, 5.5, 0.5, 22, AMBER, True)
    _add_text(slide, "Created solely to demonstrate PowerPoint RAG and cited narration.", 0.82, 4.85, 9.5, 0.4, 14, WHITE)
    _add_footer(slide, data)

    slide = prs.slides.add_slide(blank)
    _add_header(slide, "Executive summary", data["quarter"])
    _add_bullets(slide, data["summary"])
    _add_footer(slide, data)

    slide = prs.slides.add_slide(blank)
    _add_header(slide, "Quarterly KPI scorecard", "Governed finance and insurance measures")
    _add_kpi_grid(slide, data)
    _add_text(
        slide,
        f"Severity {data['severity']} · Approval {data['approval']} · "
        f"Settlement {data['settlement']} · Expense ratio {data['expense_ratio']}",
        0.75,
        6.25,
        11.8,
        0.35,
        13,
        WHITE,
    )
    _add_footer(slide, data)

    slide = prs.slides.add_slide(blank)
    _add_header(slide, "Performance by line of business", data["quarter"])
    _add_table(slide, data["lob"])
    _add_footer(slide, data)

    slide = prs.slides.add_slide(blank)
    _add_header(slide, "Regional claims performance", data["quarter"])
    _add_table(slide, data["regions"])
    _add_footer(slide, data)

    slide = prs.slides.add_slide(blank)
    _add_header(slide, "Claims and underwriting drivers", data["quarter"])
    _add_bullets(slide, data["drivers"])
    _add_footer(slide, data)

    slide = prs.slides.add_slide(blank)
    _add_header(slide, "Capital, liquidity and outlook", data["quarter"])
    _add_text(slide, "Solvency coverage", 0.9, 1.6, 3.0, 0.35, 14, MUTED)
    _add_text(slide, data["capital"], 0.9, 2.0, 3.0, 0.7, 30, GREEN, True)
    _add_text(slide, "Cash and liquid investments", 4.4, 1.6, 4.0, 0.35, 14, MUTED)
    _add_text(slide, data["cash"], 4.4, 2.0, 3.5, 0.7, 30, CYAN, True)
    _add_text(slide, "Management actions", 0.9, 3.05, 4.5, 0.4, 18, WHITE, True)
    _add_bullets(slide, data["actions"], left=1.0, top=3.55, width=11.2, height=2.8)
    _add_footer(slide, data)

    output_path = OUT_DIR / data["filename"]
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    for quarter in (Q1, Q2):
        print(build_deck(quarter))
