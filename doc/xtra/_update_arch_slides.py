"""Targeted minor text updates to slides 4-6 (architecture alignment)."""
import shutil
from pptx import Presentation

SRC = r"E:\ai_data_rag\ai_data_v2\doc\Capgemini_AI_Data_Platform_V11.pptx"
BAK = r"E:\ai_data_rag\ai_data_v2\doc\Capgemini_AI_Data_Platform_V11.backup.pptx"

# shape_index -> new_text (only when current text matches expected)
SLIDE4 = {
    1: "ASK-DB · AI Data Copilot",
    20: "Semantic\nContext",
}
SLIDE4_EXPECT = {
    1: "AI Data Copilot - AskData",
    20: "Query \nPlanning",
}

SLIDE5 = {
    3: "ASK-DB · AI Data Copilot",
    9: "Industry Pack\nSwitch",
    11: "YAML Semantic\nModel Load",
    12: "Semantic Context\nBuilder",
    17: "Follow-up &\nSQL Anchor",
    26: "Auto-Join &\nWorking Dataset",
    29: "Glossary & Metric\nLookup",
}
SLIDE5_EXPECT = {
    3: "AI Data Copilot - AskData",
    9: "Model Update",
    11: "Business Semantic Model Creation",
    12: "SQL Generation + Metadata Enrichment",
    17: "SQL Generation + Metadata Enrichment",
    26: "Semantic Model Update",
    29: "SQL Generation + Metadata Enrichment",
}

SLIDE6 = {
    1: "ASK-DB · AI Data Copilot",
    14: "CSV Multi-table\nData Upload",
    15: "CSV files · semantic join",
    19: "Industry Pack\nLoader",
    20: "YAML domain packs",
    87: "Chat Room · Trust · Narration",
    99: "Upload → Join → Semantic → SQL → Insight",
}
SLIDE6_EXPECT = {
    1: "AI Data Copilot - AskData",
    14: "Multi-format\nData Upload",
    15: "CSV · Excel · JSON · Parquet",
    19: "Config JSON\nLoader",
    20: "Industry/domain config",
    87: "Summary · Facts · Reco",
    99: "User Intent → Semantic Layer → SQL Engine → Insight",
}


def set_shape_text(slide, idx, expected, new_text, where):
    shape = slide.shapes[idx]
    if not shape.has_text_frame:
        raise ValueError(f"{where}: shape {idx} has no text")
    current = shape.text_frame.text
    if current != expected:
        raise ValueError(
            f"{where}: shape {idx} text mismatch.\n  expected: {expected!r}\n  actual:   {current!r}"
        )
    shape.text_frame.text = new_text
    return current, new_text


def main():
    shutil.copy2(SRC, BAK)
    prs = Presentation(SRC)
    changes = []

    for idx, new in SLIDE4.items():
        old, nw = set_shape_text(prs.slides[3], idx, SLIDE4_EXPECT[idx], new, "slide4")
        changes.append(("slide4", idx, old, nw))

    for idx, new in SLIDE5.items():
        old, nw = set_shape_text(prs.slides[4], idx, SLIDE5_EXPECT[idx], new, "slide5")
        changes.append(("slide5", idx, old, nw))

    for idx, new in SLIDE6.items():
        old, nw = set_shape_text(prs.slides[5], idx, SLIDE6_EXPECT[idx], new, "slide6")
        changes.append(("slide6", idx, old, nw))

    prs.save(SRC)
    print(f"Updated {len(changes)} labels in {SRC}")
    for where, idx, old, new in changes:
        print(f"  [{where}#{idx}] {old.replace(chr(10),' | ')} => {new.replace(chr(10),' | ')}")


if __name__ == "__main__":
    main()
