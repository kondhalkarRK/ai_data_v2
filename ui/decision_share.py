"""
ui/decision_share.py
Executive brief formatting, share exports (copy / email / Teams / PDF / PPT), and pin storage.
Uses mailto + downloads — no OAuth required (enterprise Graph/Teams API can be added later).
"""
from __future__ import annotations

import html
import io
import json
import urllib.parse
from datetime import datetime
from typing import Any

import pandas as pd
import streamlit as st

try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    _MPL_OK = True
except ImportError:
    _MPL_OK = False

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    _PDF_OK = True
except ImportError:
    _PDF_OK = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    _PPT_OK = True
except ImportError:
    _PPT_OK = False


def _now_label() -> str:
    return datetime.now().strftime("%d %b %Y, %H:%M")


def _evidence_line(evidence: dict | None) -> str:
    ev = evidence or {}
    parts = []
    path = ev.get("execution_path")
    if path:
        parts.append(f"Path: {path}")
    if ev.get("resolution_source"):
        parts.append(f"Source: {ev['resolution_source']}")
    trust = ev.get("trust_score")
    if trust is not None:
        parts.append(f"Trust: {trust}/100")
    return " · ".join(parts) if parts else "ASK - DB — governed semantic analytics"


def build_share_payload(
    *,
    question: str = "",
    narration: dict | None = None,
    result_df: pd.DataFrame | None = None,
    evidence: dict | None = None,
    elapsed: float | None = None,
    sql: str | None = None,
    chart_x: str | None = None,
    chart_y: str | None = None,
    chart_type: str | None = None,
) -> dict[str, Any]:
    """Normalize query/chat bundle into a shareable decision payload."""
    narr = narration or {}
    headline = str(narr.get("headline") or "Decision brief")
    raw = str(narr.get("narrative_text") or narr.get("summary") or "")
    paras = [p.strip() for p in raw.split("\n\n") if p.strip()]
    if not paras and raw.strip():
        paras = [raw.strip()]
    findings = [str(f).strip() for f in (narr.get("key_findings") or []) if str(f).strip()]
    rec = str(narr.get("recommendation") or "").strip()

    row_count = len(result_df) if isinstance(result_df, pd.DataFrame) else 0
    col_count = result_df.shape[1] if isinstance(result_df, pd.DataFrame) else 0

    return {
        "question": (question or "").strip(),
        "headline": headline,
        "paragraphs": paras,
        "findings": findings,
        "recommendation": rec,
        "result_df": result_df,
        "evidence": evidence or {},
        "elapsed": elapsed,
        "sql": (sql or "").strip(),
        "generated_at": _now_label(),
        "row_count": row_count,
        "col_count": col_count,
        "chart_x": chart_x,
        "chart_y": chart_y,
        "chart_type": chart_type,
    }


def format_executive_brief(payload: dict[str, Any], *, teams: bool = False) -> str:
    """Plain-text executive brief for copy, email, or Teams paste."""
    q = payload.get("question") or "Analytics insight"
    lines = [
        "═" * 52,
        "DECISION BRIEF — ASK - DB",
        f"Generated: {payload.get('generated_at', _now_label())}",
        "═" * 52,
        "",
        f"Question: {q}",
        "",
        f"▶ {payload.get('headline', 'Insight')}",
        "",
    ]
    for p in payload.get("paragraphs") or []:
        lines.append(p)
        lines.append("")
    for i, f in enumerate(payload.get("findings") or [], 1):
        lines.append(f"  {i}. {f}")
    if payload.get("findings"):
        lines.append("")
    if payload.get("recommendation"):
        lines.append(f"Recommendation: {payload['recommendation']}")
        lines.append("")
    rc = payload.get("row_count", 0)
    if rc:
        lines.append(f"Data: {rc:,} rows × {payload.get('col_count', 0)} columns")
    if payload.get("elapsed") is not None:
        lines.append(f"Query time: {payload['elapsed']}s")
    lines.append(_evidence_line(payload.get("evidence")))
    lines.append("")
    lines.append("— Shared from Chat Room · ASK - DB")
    if teams:
        lines.insert(0, "**Decision Brief** (paste into Microsoft Teams)")
    return "\n".join(lines).strip()


def format_teams_message(payload: dict[str, Any]) -> str:
    """Markdown-friendly block for Teams chat paste."""
    brief = format_executive_brief(payload)
    return (
        "**📊 Decision Brief**\n\n"
        f"**{payload.get('headline', 'Insight')}**\n\n"
        + "\n".join(f"• {p}" for p in (payload.get("paragraphs") or [])[:4])
        + "\n\n"
        + (f"**Recommendation:** {payload['recommendation']}\n\n" if payload.get("recommendation") else "")
        + f"_Question: {payload.get('question', '')}_\n"
        + f"_{_evidence_line(payload.get('evidence'))}_"
    )


def _resolve_chart_axes(
    result_df: pd.DataFrame,
    question: str = "",
    chart_x: str | None = None,
    chart_y: str | None = None,
    chart_type: str | None = None,
) -> tuple[str, str, str]:
    """Pick X/Y columns and chart type — mirrors Query tab chart defaults."""
    cols = list(result_df.columns)
    nums = result_df.select_dtypes(include="number").columns.tolist()
    strs = result_df.select_dtypes(exclude="number").columns.tolist()
    ct = (chart_type or "").strip().capitalize()
    if ct not in ("Bar", "Line", "Pie", "Scatter", "Area"):
        try:
            from core.chart_engine import auto_chart_type
            ct = auto_chart_type(result_df, question or "")
        except Exception:
            ct = "Bar"
    x_col = chart_x if chart_x in cols else (strs[0] if strs else cols[0])
    y_col = chart_y if chart_y in cols else (nums[0] if nums else cols[-1])
    if x_col == y_col and len(cols) > 1:
        y_col = nums[0] if nums and nums[0] != x_col else cols[-1]
    return x_col, y_col, ct


def _chart_png_bytes(
    result_df: pd.DataFrame | None,
    question: str = "",
    *,
    chart_x: str | None = None,
    chart_y: str | None = None,
    chart_type: str | None = None,
) -> bytes | None:
    if not _MPL_OK or result_df is None or result_df.empty:
        return None
    try:
        nums = result_df.select_dtypes(include="number").columns.tolist()
        if not nums:
            return None
        plot_df = result_df.head(20).copy()
        x_col, y_col, ct = _resolve_chart_axes(
            plot_df, question, chart_x, chart_y, chart_type,
        )
        if x_col not in plot_df.columns or y_col not in plot_df.columns:
            return None
        y_vals = pd.to_numeric(plot_df[y_col], errors="coerce")
        if y_vals.isna().all():
            return None
        x_vals = plot_df[x_col].astype(str).tolist()
        y_plot = y_vals.fillna(0).astype(float).tolist()

        fig, ax = plt.subplots(figsize=(7.5, 3.4))
        fig.patch.set_facecolor("white")
        ax.set_facecolor("#fafafa")
        title = str(question or "Chart")[:72]
        color = "#6366f1"
        if ct == "Line":
            ax.plot(range(len(x_vals)), y_plot, color=color, marker="o", linewidth=2, markersize=4)
            ax.set_xticks(range(len(x_vals)))
            ax.set_xticklabels(x_vals, rotation=35, ha="right", fontsize=7)
        elif ct == "Pie" and len(x_vals) <= 12:
            ax.pie(y_plot, labels=x_vals, autopct="%1.0f%%", textprops={"fontsize": 7})
        elif ct == "Area":
            ax.fill_between(range(len(x_vals)), y_plot, alpha=0.35, color=color)
            ax.plot(range(len(x_vals)), y_plot, color=color, linewidth=2)
            ax.set_xticks(range(len(x_vals)))
            ax.set_xticklabels(x_vals, rotation=35, ha="right", fontsize=7)
        elif ct == "Scatter":
            ax.scatter(range(len(x_vals)), y_plot, color=color, s=36)
            ax.set_xticks(range(len(x_vals)))
            ax.set_xticklabels(x_vals, rotation=35, ha="right", fontsize=7)
        else:
            ax.bar(range(len(x_vals)), y_plot, color=color, width=0.65)
            ax.set_xticks(range(len(x_vals)))
            ax.set_xticklabels(x_vals, rotation=35, ha="right", fontsize=7)
        ax.set_title(title, fontsize=10, color="#1e293b", pad=8)
        ax.tick_params(axis="y", labelsize=7)
        ax.grid(axis="y", alpha=0.25, linestyle="--")
        fig.tight_layout()
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=130, bbox_inches="tight", facecolor="white")
        plt.close(fig)
        buf.seek(0)
        return buf.read()
    except Exception:
        return None


def _chart_export_kwargs(payload: dict[str, Any]) -> dict[str, Any]:
    return {
        "chart_x": payload.get("chart_x"),
        "chart_y": payload.get("chart_y"),
        "chart_type": payload.get("chart_type"),
    }


def payload_headline(_df) -> str:
    return "Chart"


def build_html_brief(payload: dict[str, Any]) -> str:
    """Print-friendly HTML brief (PDF fallback) with chart + table."""
    import base64

    paras = "".join(
        f"<p>{html.escape(p)}</p>" for p in (payload.get("paragraphs") or [])
    )
    findings = "".join(
        f"<li>{html.escape(f)}</li>" for f in (payload.get("findings") or [])
    )
    rec = html.escape(payload.get("recommendation") or "")

    chart_html = ""
    ck = _chart_export_kwargs(payload)
    png = _chart_png_bytes(
        payload.get("result_df"), payload.get("question", ""), **ck,
    )
    if png:
        b64 = base64.b64encode(png).decode("ascii")
        chart_html = f'<h2>Chart</h2><img src="data:image/png;base64,{b64}" style="max-width:100%;height:auto;" alt="Chart"/>'

    table_html = ""
    table_data = _table_rows_for_export(payload.get("result_df"))
    if table_data:
        header = "".join(f"<th>{html.escape(c)}</th>" for c in table_data[0])
        body_rows = ""
        for row in table_data[1:]:
            body_rows += "<tr>" + "".join(f"<td>{html.escape(c)}</td>" for c in row) + "</tr>"
        table_html = (
            "<h2>Data table</h2>"
            f'<table border="1" cellpadding="6" cellspacing="0" style="border-collapse:collapse;width:100%;font-size:12px">'
            f"<thead><tr>{header}</tr></thead><tbody>{body_rows}</tbody></table>"
        )

    return f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>Decision Brief</title>
<style>
body{{font-family:Segoe UI,Arial,sans-serif;max-width:720px;margin:40px auto;padding:0 24px;color:#1e293b}}
h1{{font-size:22px;color:#4338ca}} h2{{font-size:16px;color:#4338ca;margin-top:24px}}
.meta{{color:#64748b;font-size:13px}}
.rec{{background:#f1f5f9;padding:12px;border-radius:8px;margin-top:16px}}
table th{{background:#eef2ff;text-align:left}}
footer{{margin-top:32px;font-size:12px;color:#94a3b8;border-top:1px solid #e2e8f0;padding-top:12px}}
</style></head><body>
<h1>{html.escape(payload.get('headline', 'Decision Brief'))}</h1>
<p class="meta">Question: {html.escape(payload.get('question') or '')} · {html.escape(payload.get('generated_at', ''))}</p>
{paras}
{"<ul>" + findings + "</ul>" if findings else ""}
{"<div class='rec'><strong>Recommendation:</strong> " + rec + "</div>" if rec else ""}
{chart_html}
{table_html}
<footer>ASK - DB · Chat Room · {_evidence_line(payload.get('evidence'))}</footer>
</body></html>"""


def _table_rows_for_export(result_df: pd.DataFrame | None, max_rows: int = 15, max_cols: int = 6) -> list[list[str]]:
    if result_df is None or result_df.empty:
        return []
    df = result_df.head(max_rows).copy()
    cols = list(df.columns[:max_cols])
    rows = [cols]
    for _, row in df.iterrows():
        rows.append([str(row[c])[:28] for c in cols])
    return rows


def build_pdf_bytes(payload: dict[str, Any]) -> bytes | None:
    if not _PDF_OK:
        return None
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import cm
        from reportlab.platypus import (
            Image as RLImage,
            Paragraph,
            SimpleDocTemplate,
            Spacer,
            Table,
            TableStyle,
        )

        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4, leftMargin=1.5 * cm, rightMargin=1.5 * cm,
                                topMargin=1.5 * cm, bottomMargin=1.5 * cm)
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle("Title", parent=styles["Heading1"], fontSize=16, spaceAfter=8)
        body_style = ParagraphStyle("Body", parent=styles["Normal"], fontSize=10, leading=14)
        meta_style = ParagraphStyle("Meta", parent=styles["Normal"], fontSize=9, textColor=colors.grey)

        story = []
        story.append(Paragraph(payload.get("headline") or "Decision Brief", title_style))
        story.append(Paragraph(
            f"Question: {html.escape(payload.get('question') or '')}<br/>"
            f"Generated: {html.escape(payload.get('generated_at', ''))}",
            meta_style,
        ))
        story.append(Spacer(1, 0.3 * cm))

        for para in (payload.get("paragraphs") or [])[:4]:
            story.append(Paragraph(html.escape(para), body_style))
            story.append(Spacer(1, 0.15 * cm))
        if payload.get("recommendation"):
            story.append(Paragraph(
                f"<b>Recommendation:</b> {html.escape(payload['recommendation'])}",
                body_style,
            ))

        ck = _chart_export_kwargs(payload)
        png = _chart_png_bytes(
            payload.get("result_df"), payload.get("question", ""), **ck,
        )
        table_data = _table_rows_for_export(payload.get("result_df"))

        if png:
            story.append(Spacer(1, 0.45 * cm))
            story.append(Paragraph("Chart", styles["Heading2"]))
            story.append(Spacer(1, 0.15 * cm))
            story.append(RLImage(io.BytesIO(png), width=16 * cm, height=6.5 * cm))
        if table_data:
            story.append(Spacer(1, 0.35 * cm))
            story.append(Paragraph("Data table", styles["Heading2"]))
            story.append(Spacer(1, 0.15 * cm))
            tbl = Table(table_data, repeatRows=1)
            tbl.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#eef2ff")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#1e293b")),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
                ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#cbd5e1")),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f8fafc")]),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]))
            story.append(tbl)

        story.append(Spacer(1, 0.5 * cm))
        story.append(Paragraph(html.escape(_evidence_line(payload.get("evidence"))[:120]), meta_style))

        doc.build(story)
        buf.seek(0)
        return buf.read()
    except Exception:
        return None


def _wrap_text(text: str, width: int) -> list[str]:
    words = str(text).split()
    lines, cur = [], ""
    for w in words:
        if len(cur) + len(w) + 1 <= width:
            cur = (cur + " " + w).strip()
        else:
            if cur:
                lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines or [""]


def build_pptx_bytes(payload: dict[str, Any]) -> bytes | None:
    if not _PPT_OK:
        return None
    try:
        from pptx.util import Inches, Pt

        prs = Presentation()
        # Slide 1 — brief
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide1.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(1))
        tb.text_frame.text = payload.get("headline") or "Decision Brief"
        tb.text_frame.paragraphs[0].font.size = Pt(24)
        tb.text_frame.paragraphs[0].font.bold = True
        sub = slide1.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(9), Inches(0.45))
        sub.text_frame.text = f"{payload.get('question', '')} · {payload.get('generated_at', '')}"
        sub.text_frame.paragraphs[0].font.size = Pt(11)
        body = slide1.shapes.add_textbox(Inches(0.5), Inches(1.75), Inches(9), Inches(4.8))
        btf = body.text_frame
        btf.word_wrap = True
        first = True
        for para in (payload.get("paragraphs") or [])[:3]:
            p = btf.paragraphs[0] if first else btf.add_paragraph()
            first = False
            p.text = para
            p.font.size = Pt(13)
        if payload.get("recommendation"):
            p = btf.add_paragraph()
            p.text = f"Recommendation: {payload['recommendation']}"
            p.font.size = Pt(12)
            p.font.bold = True

        rdf = payload.get("result_df")
        ck = _chart_export_kwargs(payload)
        png = _chart_png_bytes(rdf, payload.get("question", ""), **ck)
        table_data = _table_rows_for_export(rdf)

        if png or table_data:
            slide2 = prs.slides.add_slide(prs.slide_layouts[6])
            y_off = 0.35
            if png:
                t2 = slide2.shapes.add_textbox(Inches(0.5), Inches(y_off), Inches(9), Inches(0.45))
                t2.text_frame.text = "Chart"
                t2.text_frame.paragraphs[0].font.size = Pt(16)
                t2.text_frame.paragraphs[0].font.bold = True
                slide2.shapes.add_picture(
                    io.BytesIO(png), Inches(0.55), Inches(y_off + 0.55), width=Inches(8.9),
                )
                y_off = 4.15
            if table_data:
                t3 = slide2.shapes.add_textbox(Inches(0.5), Inches(y_off), Inches(9), Inches(0.45))
                t3.text_frame.text = "Data table"
                t3.text_frame.paragraphs[0].font.size = Pt(16)
                t3.text_frame.paragraphs[0].font.bold = True
                nrows = len(table_data)
                ncols = len(table_data[0]) if table_data else 1
                tbl_h = min(0.32 * nrows + 0.25, 2.8 if png else 5.5)
                tbl_shape = slide2.shapes.add_table(
                    nrows, ncols,
                    Inches(0.45), Inches(y_off + 0.5),
                    Inches(9.1), Inches(tbl_h),
                )
                table = tbl_shape.table
                for ri, row in enumerate(table_data):
                    for ci, val in enumerate(row):
                        cell = table.cell(ri, ci)
                        cell.text = str(val)
                        if ri == 0:
                            cell.text_frame.paragraphs[0].font.bold = True
                            cell.text_frame.paragraphs[0].font.size = Pt(9)
                        else:
                            cell.text_frame.paragraphs[0].font.size = Pt(8)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()
    except Exception:
        return None


def pin_decision(payload: dict[str, Any]) -> None:
    pins = list(st.session_state.get("pinned_decisions") or [])
    pin_id = f"pin_{datetime.now().strftime('%Y%m%d%H%M%S')}_{len(pins)}"
    sample = None
    rdf = payload.get("result_df")
    if isinstance(rdf, pd.DataFrame) and not rdf.empty:
        sample = rdf.head(5).to_csv(index=False)
    entry = {
        "id": pin_id,
        "question": payload.get("question", ""),
        "headline": payload.get("headline", ""),
        "summary": " ".join((payload.get("paragraphs") or [""])[0].split()[:40]),
        "recommendation": payload.get("recommendation", ""),
        "pinned_at": _now_label(),
        "row_count": payload.get("row_count", 0),
        "sample_csv": sample,
        "suggested_question": payload.get("question", ""),
    }
    pins.insert(0, entry)
    st.session_state["pinned_decisions"] = pins[:8]


def unpin_decision(pin_id: str) -> None:
    pins = [p for p in (st.session_state.get("pinned_decisions") or []) if p.get("id") != pin_id]
    st.session_state["pinned_decisions"] = pins


def render_pinned_strip(on_ask) -> None:
    """Compact pinned-decisions popover shown above the Chat Room thread."""
    pins = st.session_state.get("pinned_decisions") or []
    if not pins:
        return
    with st.popover(f"📌 Pinned ({len(pins)})", help="Open or remove pinned decisions"):
        for pin in pins[:8]:
            st.markdown(
                f'<div class="dr-pin-card dr-pin-card-compact"><div class="dr-pin-headline">'
                f'{html.escape(str(pin.get("headline", "Insight"))[:48])}</div>'
                f'<div class="dr-pin-meta">{html.escape(str(pin.get("pinned_at", "")))}'
                f' · {pin.get("row_count", 0):,} rows</div></div>',
                unsafe_allow_html=True,
            )
            c1, c2 = st.columns([4, 1])
            with c1:
                if st.button("Open", key=f"pin_open_{pin['id']}", use_container_width=True):
                    q = pin.get("suggested_question") or pin.get("question")
                    if q and callable(on_ask):
                        on_ask(q)
                        st.rerun()
            with c2:
                if st.button("✕", key=f"pin_rm_{pin['id']}", use_container_width=True):
                    unpin_decision(pin["id"])
                    st.rerun()


def render_proactive_landing(working_df: pd.DataFrame, insights: list[dict], on_ask) -> None:
    """Brief-style empty state for Chat Room."""
    st.markdown(
        """
        <div class="dr-landing">
          <div class="dr-landing-eyebrow">CHAT ROOM</div>
          <div class="dr-landing-title">Today's data priorities</div>
          <div class="dr-landing-sub">Insights from your dataset — ask, refine, pin, or share.</div>
        </div>
        """,
        unsafe_allow_html=True,
    )
    if not insights:
        st.info("Load data and run a query — proactive priorities will appear here.")
        return
    st.markdown('<div class="dr-priority-grid">', unsafe_allow_html=True)
    cols = st.columns(min(len(insights), 3))
    for i, ins in enumerate(insights[:3]):
        with cols[i]:
            direction = ins.get("direction", "neutral")
            icon = {"up": "🟢", "down": "🔴", "neutral": "🟡"}.get(direction, "🔵")
            st.markdown(
                f'<div class="dr-priority-card">'
                f'<div class="dr-priority-icon">{icon}</div>'
                f'<div class="dr-priority-title">{html.escape(str(ins.get("title", "")))}</div>'
                f'<div class="dr-priority-summary">{html.escape(str(ins.get("summary", "")))}</div>'
                f'</div>',
                unsafe_allow_html=True,
            )
            sq = ins.get("suggested_question")
            if sq and st.button("Explore →", key=f"dr_proactive_{i}", use_container_width=True):
                if callable(on_ask):
                    on_ask(sq)
                    st.rerun()
    st.markdown("</div>", unsafe_allow_html=True)


def _clipboard_js(text: str) -> str:
    escaped = json.dumps(text)
    return f"""
    <script>
    navigator.clipboard.writeText({escaped}).then(function() {{
        window.parent.document.body.dispatchEvent(new CustomEvent('decision-copied'));
    }});
    </script>
    """


def render_share_and_pin(
    payload: dict[str, Any],
    key_prefix: str,
    *,
    show_pin: bool = True,
) -> None:
    """Compact icon-only Pin + Share for query/chat results."""
    if not payload.get("headline") and not payload.get("question"):
        return

    st.markdown('<div class="dr-icon-actions">', unsafe_allow_html=True)
    ic1, ic2 = st.columns([1, 1], gap="small")
    with ic1:
        if show_pin and st.button(
            "📌",
            key=f"{key_prefix}_pin",
            help="Pin to Chat Room",
        ):
            pin_decision(payload)
            st.toast("Pinned to Chat Room", icon="📌")
            st.rerun()
    with ic2:
        with st.popover("📤", help="Share decision brief"):
            st.caption("Share brief, chart, and table.")
            brief = format_executive_brief(payload)
            teams_msg = format_teams_message(payload)

            if st.button("📋 Copy brief", key=f"{key_prefix}_copy", use_container_width=True):
                st.session_state[f"{key_prefix}_copy_text"] = brief
                st.toast("Brief ready — copy from box below", icon="📋")

            copy_text = st.session_state.get(f"{key_prefix}_copy_text")
            if copy_text:
                st.text_area("Copy to clipboard", value=copy_text, height=100, key=f"{key_prefix}_ta")

            subject = urllib.parse.quote(f"Decision Brief: {(payload.get('headline') or 'Insight')[:60]}")
            body = urllib.parse.quote(brief[:3500])
            st.markdown(
                f'<a href="mailto:?subject={subject}&body={body}" target="_blank" '
                f'class="dr-share-link">📧 Open in Outlook / Email</a>',
                unsafe_allow_html=True,
            )

            if st.button("💬 Copy for Teams", key=f"{key_prefix}_teams", use_container_width=True):
                st.session_state[f"{key_prefix}_teams_text"] = teams_msg
                st.toast("Teams message ready — copy below", icon="💬")

            teams_text = st.session_state.get(f"{key_prefix}_teams_text")
            if teams_text:
                st.text_area("Paste into Teams", value=teams_text, height=90, key=f"{key_prefix}_teams_ta")

            # PDF/PPT chart generation is expensive. Streamlit executes closed
            # popover bodies on every rerun, so prepare files only on request.
            export_key = f"{key_prefix}_prepared_exports"
            if st.button(
                "Prepare PDF / PowerPoint",
                key=f"{key_prefix}_prepare_exports",
                use_container_width=True,
            ):
                with st.spinner("Preparing export files…"):
                    st.session_state[export_key] = {
                        "pdf": build_pdf_bytes(payload),
                        "ppt": build_pptx_bytes(payload),
                    }

            exports = st.session_state.get(export_key) or {}
            pdf_bytes = exports.get("pdf")
            if pdf_bytes:
                st.download_button(
                    "📄 Download PDF (brief + chart + table)",
                    data=pdf_bytes,
                    file_name="decision_brief.pdf",
                    mime="application/pdf",
                    key=f"{key_prefix}_pdf",
                    use_container_width=True,
                )
            elif export_key in st.session_state:
                st.download_button(
                    "📄 Download HTML brief",
                    data=build_html_brief(payload).encode("utf-8"),
                    file_name="decision_brief.html",
                    mime="text/html",
                    key=f"{key_prefix}_html",
                    use_container_width=True,
                )

            ppt_bytes = exports.get("ppt")
            if ppt_bytes:
                st.download_button(
                    "📊 Download PowerPoint (brief + chart + table)",
                    data=ppt_bytes,
                    file_name="decision_brief.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key=f"{key_prefix}_ppt",
                    use_container_width=True,
                )
            elif export_key in st.session_state:
                st.caption("Install python-pptx for PowerPoint export.")

            rdf = payload.get("result_df")
            if isinstance(rdf, pd.DataFrame) and not rdf.empty:
                st.download_button(
                    "⬇️ Download data (CSV)",
                    data=rdf.to_csv(index=False).encode(),
                    file_name="decision_data.csv",
                    mime="text/csv",
                    key=f"{key_prefix}_csv",
                    use_container_width=True,
                )
    st.markdown("</div>", unsafe_allow_html=True)
